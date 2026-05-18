from __future__ import annotations

from pathlib import Path

import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.analytics import aggregate_weekly_sales, prepare_weekly_sales_frame, weekly_kpi_summary
from src.config import RAW_DATA_PATH, REPORT_DIR, SCREENSHOT_DIR


FALLBACK_ASSET_DIR = Path("C:/tmp")


def _safe_image_save(image: Image.Image, path: Path) -> Path:
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        image.save(path)
        return path
    except OSError:
        fallback = FALLBACK_ASSET_DIR / path.name
        fallback.parent.mkdir(parents=True, exist_ok=True)
        image.save(fallback)
        return fallback


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/seguisb.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def _draw_metric(draw: ImageDraw.ImageDraw, x: int, y: int, label: str, value: str) -> None:
    draw.rounded_rectangle((x, y, x + 270, y + 120), radius=18, fill="#F7FAFC", outline="#D8E2EA", width=2)
    draw.text((x + 22, y + 20), label.upper(), fill="#5B6770", font=_font(18, True))
    draw.text((x + 22, y + 54), value, fill="#0F172A", font=_font(34, True))


def _bar_chart(draw: ImageDraw.ImageDraw, data: pd.DataFrame, x: int, y: int, width: int, height: int) -> None:
    values = data["sales"].tolist()
    labels = [f"Store {int(store)}" for store in data["store"].tolist()]
    max_value = max(values)
    bar_gap = 16
    bar_height = int((height - bar_gap * (len(values) - 1)) / len(values))
    for index, (label, value) in enumerate(zip(labels, values)):
        top = y + index * (bar_height + bar_gap)
        draw.text((x, top + 8), label, fill="#27313A", font=_font(18, True))
        bar_x = x + 115
        bar_width = int((value / max_value) * (width - 130))
        draw.rounded_rectangle((bar_x, top, bar_x + bar_width, top + bar_height), radius=9, fill="#137C8B")
        draw.text((bar_x + bar_width + 10, top + 8), f"${value / 1_000_000:.1f}M", fill="#27313A", font=_font(17))


def create_dashboard_screenshots() -> list[Path]:
    df = prepare_weekly_sales_frame(pd.read_csv(RAW_DATA_PATH))
    summary = weekly_kpi_summary(df)
    store_sales = aggregate_weekly_sales(df, "store").head(5)

    image = Image.new("RGB", (1440, 900), "#FFFFFF")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 1440, 92), fill="#0F172A")
    draw.text((54, 26), "Walmart Sales Analytics Dashboard", fill="#FFFFFF", font=_font(36, True))
    draw.text((54, 112), "Real weekly sales dataset overview", fill="#0F172A", font=_font(30, True))

    _draw_metric(draw, 54, 170, "Total sales", f"${summary['total_sales'] / 1_000_000:.1f}M")
    _draw_metric(draw, 346, 170, "Stores", f"{summary['stores']}")
    _draw_metric(draw, 638, 170, "Weeks", f"{summary['weeks']}")
    _draw_metric(draw, 930, 170, "Avg weekly", f"${summary['average_weekly_sales'] / 1_000_000:.2f}M")

    draw.text((54, 350), "Top stores by sales", fill="#0F172A", font=_font(28, True))
    _bar_chart(draw, store_sales, 54, 410, 760, 300)

    draw.rounded_rectangle((890, 355, 1326, 725), radius=18, fill="#F7FAFC", outline="#D8E2EA", width=2)
    draw.text((928, 394), "Dashboard Pages", fill="#0F172A", font=_font(26, True))
    pages = ["Overview trend analysis", "Store ranking and consistency", "Holiday vs regular weeks", "Economic signal comparison", "Database or CSV data source"]
    for index, page in enumerate(pages):
        draw.ellipse((930, 455 + index * 48, 946, 471 + index * 48), fill="#137C8B")
        draw.text((964, 446 + index * 48), page, fill="#27313A", font=_font(22))

    overview_path = SCREENSHOT_DIR / "dashboard_overview.png"
    overview_path = _safe_image_save(image, overview_path)

    architecture = Image.new("RGB", (1440, 900), "#FFFFFF")
    draw = ImageDraw.Draw(architecture)
    draw.text((54, 46), "Project Architecture", fill="#0F172A", font=_font(42, True))
    boxes = [
        (80, 190, "Raw CSV", "Real Walmart dataset"),
        (390, 190, "ETL", "Pandas validation and loading"),
        (700, 190, "SQL Database", "PostgreSQL / MySQL"),
        (1010, 190, "Dashboard", "Streamlit analytics"),
        (700, 500, "SQL Views", "Reusable KPIs and queries"),
    ]
    for x, y, title, body in boxes:
        draw.rounded_rectangle((x, y, x + 250, y + 130), radius=18, fill="#F7FAFC", outline="#137C8B", width=3)
        draw.text((x + 26, y + 28), title, fill="#0F172A", font=_font(26, True))
        draw.text((x + 26, y + 72), body, fill="#4B5563", font=_font(18))
    for start, end in [((330, 255), (390, 255)), ((640, 255), (700, 255)), ((950, 255), (1010, 255)), ((825, 320), (825, 500)), ((950, 565), (1080, 320))]:
        draw.line((*start, *end), fill="#0F172A", width=4)
    architecture_path = SCREENSHOT_DIR / "project_architecture.png"
    architecture_path = _safe_image_save(architecture, architecture_path)
    return [overview_path, architecture_path]


def _add_title(slide, text: str, top: float = 0.55) -> None:
    box = slide.shapes.add_textbox(Inches(0.65), Inches(top), Inches(11.0), Inches(0.55))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(15, 23, 42)


def _add_bullets(slide, bullets: list[str], top: float = 1.55) -> None:
    box = slide.shapes.add_textbox(Inches(0.85), Inches(top), Inches(11.0), Inches(4.5))
    text_frame = box.text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = RGBColor(39, 49, 58)
        paragraph.space_after = Pt(12)


def create_presentation() -> Path:
    df = prepare_weekly_sales_frame(pd.read_csv(RAW_DATA_PATH))
    summary = weekly_kpi_summary(df)
    store_sales = aggregate_weekly_sales(df, "store").head(3)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(15, 23, 42)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(10.5), Inches(1.7))
    paragraph = title.text_frame.paragraphs[0]
    paragraph.text = "Walmart Sales Analytics Dashboard"
    paragraph.font.size = Pt(46)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    subtitle = slide.shapes.add_textbox(Inches(0.85), Inches(3.0), Inches(8.8), Inches(0.8))
    subtitle.text_frame.text = "SQL database design, ETL, real weekly sales analytics, and Streamlit reporting"
    subtitle.text_frame.paragraphs[0].font.size = Pt(22)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(209, 213, 219)

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Project Scope")
    _add_bullets(
        slide,
        [
            "Designed PostgreSQL and MySQL schemas for POS and weekly store sales analysis.",
            "Built ETL scripts for validation, transformation, and database loading.",
            "Created an interactive dashboard with CSV upload and database connection modes.",
            "Documented architecture, ER diagrams, business questions, and final findings.",
        ],
    )

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Real Dataset Snapshot")
    _add_bullets(
        slide,
        [
            f"Total sales in bundled real subset: ${summary['total_sales']:,.0f}",
            f"Stores represented: {summary['stores']}",
            f"Unique weeks represented: {summary['weeks']}",
            f"Average weekly store sales: ${summary['average_weekly_sales']:,.0f}",
        ],
    )

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Top Stores")
    rows = [[f"Store {int(row.store)}", f"${row.sales:,.0f}", f"${row.avg_weekly_sales:,.0f}"] for row in store_sales.itertuples()]
    table = slide.shapes.add_table(len(rows) + 1, 3, Inches(0.85), Inches(1.55), Inches(10.8), Inches(2.2)).table
    headers = ["Store", "Total Sales", "Avg Weekly Sales"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(19, 124, 139)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
    for row_index, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            table.cell(row_index, col).text = value

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Dashboard and SQL Deliverables")
    _add_bullets(
        slide,
        [
            "Dashboard pages: Overview, Stores, Economic Signals, Products, Customers, Operations, Data.",
            "SQL views: daily sales, product performance, customer segments, store performance, weekly trends.",
            "Analytics queries cover revenue trends, holiday impact, top products, customer segments, and economic signals.",
        ],
    )

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Recommended Next Enhancements")
    _add_bullets(
        slide,
        [
            "Load the full 6,435-row dataset using python -m src.download_dataset.",
            "Add forecasting models for weekly sales prediction.",
            "Deploy the dashboard and connect it to a managed PostgreSQL database.",
            "Add scheduled refresh and role-based executive/store-manager views.",
        ],
    )

    output_path = REPORT_DIR / "walmart_sales_analytics_presentation.pptx"
    try:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path
    except OSError:
        fallback = FALLBACK_ASSET_DIR / output_path.name
        fallback.parent.mkdir(parents=True, exist_ok=True)
        prs.save(fallback)
        return fallback


if __name__ == "__main__":
    screenshots = create_dashboard_screenshots()
    deck = create_presentation()
    print("Created screenshots:")
    for screenshot in screenshots:
        print(screenshot)
    print(f"Created presentation: {deck}")
